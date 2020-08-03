import PyPDF2
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import requests
from question_answer import qa_gen
import fitz
from rake_nltk import Rake
from datetime import datetime
ra = Rake() 



def summerizer(title, content):

    url = "https://aylien-text.p.rapidapi.com/summarize"

    querystring = {"title":title,"text": content}

    headers = {
        'x-rapidapi-host': "aylien-text.p.rapidapi.com",
        'x-rapidapi-key': "e4573ac34fmshc71719554be1369p1d0fcajsnc45c98c12e74"
        }

    response = requests.request("GET", url, headers=headers, params=querystring)

    response = response.json()

    return response['sentences']

def divide_chunks(l, n): 
      
    for i in range(0, len(l), n):  
        yield l[i:i + n] 



def pdf_to_ppt(filename,name):

    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    r = requests.get(filename, allow_redirects=True)

    open('temp.pdf', 'wb').write(r.content)
    doc = fitz.open('temp.pdf')
    net_content = ""
    for page in doc:
        content = page.getText()
        net_content = net_content + content
        ra.extract_keywords_from_text(content)
        ra.get_ranked_phrases()

        page_summary = summerizer(name, content)

        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = ra.get_ranked_phrases()[0]

        tf = body_shape.text_frame
        tf.text = 'Summary'

        for line in page_summary:
            p = tf.add_paragraph()
            #p.text = line
            run = p.add_run()
            run.text = line
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(16)
            #p.level = 1
    
    strr = "".join(str(datetime.utcnow()).split(" "))
    strr = strr.replace(":","")
    strr = strr.replace(".","")

    summary_name = 'static/summary'+strr+'.pptx'
    prs.save(summary_name)
    directory = os.path.dirname(__file__)

    qa_pair = qa_gen(net_content)

    return directory, summary_name, qa_pair

#d, f, qa = pdf_to_ppt('sample.pdf')
    
#print(qa)